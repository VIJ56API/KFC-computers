import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Clean Light Theme (Light Blue & Golden)
    COLOR_BG = RGBColor(248, 250, 252)       # Light Slate
    COLOR_CARD = RGBColor(255, 255, 255)     # Pure White
    COLOR_TEXT = RGBColor(0, 0, 0)           # High-Contrast Black
    COLOR_MUTED = RGBColor(30, 41, 59)       # Dark Slate
    COLOR_BLUE = RGBColor(2, 132, 199)       # Light Blue
    COLOR_GOLD = RGBColor(217, 119, 6)       # Golden
    COLOR_BORDER = RGBColor(203, 213, 225)   # Border Slate

    blank_slide_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="KFC COMPUTERS PRESENTATION"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_BLUE

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT

    def set_slide_background(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    card1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_CARD
    card1.line.color.rgb = COLOR_BORDER
    card1.line.width = Pt(2)

    tf1 = card1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "KFC COMPUTERS E-COMMERCE PLATFORM"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.alignment = PP_ALIGN.CENTER

    p2 = tf1.add_paragraph()
    p2.text = "High-Performance PCs, Custom Builder & Admin Management System"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_BLUE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf1.add_paragraph()
    p3.text = "\nFull E-Commerce Journey • Interactive Compatibility Engine • JSON Catalog CRUD"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_MUTED
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf1.add_paragraph()
    p4.text = "\nPresented by: Development Team | Platform: PHP & SQLite/JSON"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_TEXT
    p4.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "Executive Summary & System Purpose")

    points2 = [
        ("Complete Ready-To-Use Product Catalog", "Pre-loaded with 8 realistic Ready-Made Gaming/Workstation PCs & 52 components in Indian Rupee (₹) pricing."),
        ("Interactive Custom PC Builder", "12 component category selection with real-time socket, RAM, and PSU wattage compatibility checking."),
        ("No-Code Admin Product Management", "Complete Admin Panel allowing non-technical managers to add, edit, delete products and upload images without modifying PHP code."),
        ("Robust File & DB Architecture", "Products stored in JSON files (`products.json` & `components.json`), while user credentials and orders reside in SQLite/PDO database.")
    ]

    for idx, (head, desc) in enumerate(points2):
        top_pos = Inches(1.6 + idx * 1.3)
        c = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(11.733), Inches(1.1))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_BORDER
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        p_head = tf.paragraphs[0]
        p_head.text = f"• {head}"
        p_head.font.size = Pt(16)
        p_head.font.bold = True
        p_head.font.color.rgb = COLOR_GOLD

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_TEXT

    # -------------------------------------------------------------
    # SLIDE 3: Ready-Made PC Catalog
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Pre-Configured Ready-Made PC Lineup")

    pcs = [
        ("KFC Starter Gaming PC", "₹45,999", "Ryzen 5 5600 • RTX 3050 6GB • 16GB DDR4 • 512GB SSD"),
        ("KFC Gaming Pro", "₹69,999", "Ryzen 5 7600 • RTX 4060 8GB • 16GB DDR5 • 1TB SSD"),
        ("KFC Creator PC", "₹89,999", "Core i7 14700K • RTX 4060 Ti 8GB • 32GB DDR5 • 240mm AIO"),
        ("KFC Ultimate Gaming PC", "₹1,29,999", "Ryzen 7 7800X3D • RTX 4070 12GB • 32GB DDR5 • 2TB SSD"),
        ("KFC Student PC", "₹32,999", "Core i3 14100 • Integrated UHD Graphics • 16GB RAM"),
        ("KFC Office PC", "₹38,999", "Core i5 14400 • Integrated Graphics • 512GB NVMe SSD")
    ]

    for idx, (title, price, specs) in enumerate(pcs):
        row = idx // 3
        col = idx % 3
        left = Inches(0.8 + col * 3.9)
        top = Inches(1.6 + row * 2.6)

        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_BORDER
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT

        p2 = tf.add_paragraph()
        p2.text = price
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_GOLD

        p3 = tf.add_paragraph()
        p3.text = f"\n{specs}"
        p3.font.size = Pt(12)
        p3.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 4: Interactive Custom PC Builder
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Interactive Custom PC Builder & Compatibility Engine")

    col1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2))
    col1.fill.solid()
    col1.fill.fore_color.rgb = COLOR_CARD
    col1.line.color.rgb = COLOR_BORDER
    col1.line.width = Pt(1.5)

    tf1 = col1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "12 Component Categories"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE

    cats_text = """
• CPU (Processors)
• Motherboard
• RAM (Memory)
• GPU (Graphics)
• SSD (NVMe Storage)
• HDD (Mass Storage)
• PSU (Power Supply)
• Cabinet (Chassis)
• CPU Cooler
• Monitor (Displays)
• Mechanical Keyboard
• Gaming Mouse
"""
    p_cats = tf1.add_paragraph()
    p_cats.text = cats_text
    p_cats.font.size = Pt(13)
    p_cats.font.color.rgb = COLOR_TEXT

    col2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    col2.fill.solid()
    col2.fill.fore_color.rgb = COLOR_CARD
    col2.line.color.rgb = COLOR_BORDER
    col2.line.width = Pt(1.5)

    tf2 = col2.text_frame
    tf2.word_wrap = True
    p_comp = tf2.paragraphs[0]
    p_comp.text = "Real-Time Compatibility Rules"
    p_comp.font.size = Pt(18)
    p_comp.font.bold = True
    p_comp.font.color.rgb = COLOR_GOLD

    rules_text = """
1. CPU vs Motherboard Socket Match:
   - Validates AM4, AM5, and LGA1700 sockets.

2. Memory Type Validation:
   - Ensures Motherboard RAM slots match DDR4 or DDR5 RAM sticks.

3. PSU Wattage Calculation:
   - Calculates system power draw (CPU + GPU + 100W base) against PSU capacity.

4. Your Build Summary:
   - Live total price calculation in ₹ with one-click Add to Cart.
"""
    p_rules = tf2.add_paragraph()
    p_rules.text = rules_text
    p_rules.font.size = Pt(13)
    p_rules.font.color.rgb = COLOR_TEXT

    # -------------------------------------------------------------
    # SLIDE 5: Admin Panel & CRUD
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Admin Catalog Management & Image Upload System")

    admin_features = [
        ("Product Dashboard", "View all catalog items with images, prices in ₹, stock levels, categories, and direct action triggers."),
        ("Add New Product", "Interactive form for Ready-Made PCs and components with brand, price, stock, specs, and socket parameters."),
        ("Edit Product", "Instant modification of existing prices, specifications, descriptions, and stock counts without editing source code."),
        ("Secure Image Uploads", "Upload JPG, PNG, WEBP, and SVG product images saved securely in `assets/images/products/` with unique filenames.")
    ]

    for idx, (title, desc) in enumerate(admin_features):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.6 + row * 2.6)

        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_BORDER
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_GOLD

        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_TEXT

    # Save presentation
    output_path = os.path.join(os.getcwd(), "KFC_Computers_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
